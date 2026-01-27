import streamlit as st
import mysql.connector
from mysql.connector import Error
import requests
import json
import hashlib
import pandas as pd
from typing import List, Dict, Tuple, Optional
from datetime import datetime
import os
import io
import re
import shutil
import zipfile
from io import BytesIO
# ================= RAG IMPORTS =================
try:
    import PyPDF2
    import fitz
    PYMUPDF_AVAILABLE = True
except:
    PYMUPDF_AVAILABLE = False
try:
    import chromadb
    from chromadb.config import Settings
    CHROMADB_AVAILABLE = True
except:
    CHROMADB_AVAILABLE = False
try:
    from sentence_transformers import SentenceTransformer, CrossEncoder
    SEMANTIC_AVAILABLE = True
except:
    SEMANTIC_AVAILABLE = False
try:
    from docx import Document
    DOCX_AVAILABLE = True
except:
    DOCX_AVAILABLE = False
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except:
    PPTX_AVAILABLE = False
try:
    from PIL import Image
    import easyocr
    IMAGE_AVAILABLE = True
except:
    IMAGE_AVAILABLE = False
try:
    import numpy as np
    NUMPY_AVAILABLE = True
except:
    NUMPY_AVAILABLE = False
try:
    import markdown
    from bs4 import BeautifulSoup
    MARKDOWN_AVAILABLE = True
except:
    MARKDOWN_AVAILABLE = False
# ================= CONFIGURATION =================
st.set_page_config(
    page_title="🎯 Unified AI Assistant",
    page_icon="🎯",
    layout="wide"
)
# ================= API CONFIGURATION =================
API_PROVIDER = "huggingface"
try:
    if API_PROVIDER == "gemini":
        API_KEY = st.secrets["api"]["gemini_key"]
    elif API_PROVIDER == "openai":
        API_KEY = st.secrets["api"]["openai_key"]
    else:
        API_KEY = st.secrets["api"]["hf_token"]
except:
    pass
if API_PROVIDER == "huggingface":
    API_URL = "https://router.huggingface.co/v1/chat/completions"
    MODEL_NAME = "meta-llama/Llama-3.1-8B-Instruct:fastest"
    HEADERS = {
        "Authorization": f"Bearer {API_KEY}",
        "Content-Type": "application/json"
    }
# ================= MODE CONSTANTS =================
MODE_DATABASE = "database"
MODE_RAG = "rag"
# Database Query Settings
MAX_RECENT_MESSAGES = 10
MAX_SEMANTIC_MESSAGES = 5
MAX_CONTEXT_TOKENS = 3000
# RAG Settings
SEMANTIC_CHUNK_SIZE = 400
SLIDING_WINDOW_SIZE = 600
WINDOW_OVERLAP = 150
INITIAL_RETRIEVAL_K = 20
RERANK_TOP_K = 10
FINAL_CONTEXT_K = 5
# ================= CACHE MODELS =================
@st.cache_resource
def load_rag_models():
    """Load RAG models if available"""
    if not SEMANTIC_AVAILABLE or not CHROMADB_AVAILABLE:
        return None, None, None
    
    try:
        embed = SentenceTransformer("all-MiniLM-L6-v2")
        reranker = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
        client = chromadb.Client(
            Settings(
                persist_directory="./chroma_db",
                anonymized_telemetry=False
            )
        )
        return embed, reranker, client
    except Exception as e:
        st.error(f"Error loading RAG models: {e}")
        return None, None, None
@st.cache_resource
def load_db_embedding_model():
    """Load embedding model for database semantic search"""
    if not SEMANTIC_AVAILABLE:
        return None
    try:
        return SentenceTransformer('all-MiniLM-L6-v2')
    except:
        return None
@st.cache_resource
def get_ocr_reader():
    """Load OCR reader if available"""
    if not IMAGE_AVAILABLE:
        return None
    try:
        return easyocr.Reader(['en'], gpu=False)
    except:
        return None
# ================= DATABASE CONNECTIONS =================
@st.cache_resource
def get_auth_db_connection():
    """Connect to authentication database"""
    try:
        if "auth_database" in st.secrets:
            connection = mysql.connector.connect(
                host=st.secrets["auth_database"]["host"],
                port=int(st.secrets["auth_database"]["port"]),
                database=st.secrets["auth_database"]["database"],
                user=st.secrets["auth_database"]["user"],
                password=st.secrets["auth_database"]["password"],
                ssl_disabled=False,
                ssl_verify_cert=True,
                ssl_ca=st.secrets["auth_database"].get("ssl_ca", ""),
                ssl_verify_identity=True,
                connect_timeout=30
            )
        else:
            connection = mysql.connector.connect(
                host='localhost',
                database='auth_db',
                user='root',
                password='password',
                connect_timeout=10
            )
        
        if connection.is_connected():
            init_auth_tables(connection)
            return connection
    except Error as e:
        st.error(f"❌ Auth Database connection failed: {e}")
        return None
def init_auth_tables(connection):
    """Initialize authentication tables"""
    cursor = connection.cursor()
    try:
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS users (
                user_id INT AUTO_INCREMENT PRIMARY KEY,
                username VARCHAR(50) UNIQUE NOT NULL,
                password_hash VARCHAR(64) NOT NULL,
                email VARCHAR(100),
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS chats (
                chat_id INT AUTO_INCREMENT PRIMARY KEY,
                user_id INT NOT NULL,
                title VARCHAR(255) NOT NULL,
                mode VARCHAR(20) NOT NULL DEFAULT 'database',
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users(user_id) ON DELETE CASCADE
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS chat_history (
                history_id INT AUTO_INCREMENT PRIMARY KEY,
                chat_id INT NOT NULL,
                user_id INT NOT NULL,
                question TEXT,
                query_generated TEXT,
                response TEXT,
                mode VARCHAR(20),
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (chat_id) REFERENCES chats(chat_id) ON DELETE CASCADE
            )
        """)
        
        connection.commit()
    except Error as e:
        st.error(f"Error initializing auth tables: {e}")
    finally:
        cursor.close()
@st.cache_resource
def get_business_db_connection():
    """Connect to business database"""
    try:
        if "database" in st.secrets:
            connection = mysql.connector.connect(
                host=st.secrets["database"]["host"],
                port=int(st.secrets["database"]["port"]),
                database=st.secrets["database"]["database"],
                user=st.secrets["database"]["user"],
                password=st.secrets["database"]["password"],
                ssl_disabled=False,
                ssl_verify_cert=True,
                ssl_ca=st.secrets["database"]["ssl_ca"],
                ssl_verify_identity=True,
                connect_timeout=30
            )
        else:
            connection = mysql.connector.connect(
                host='localhost',
                database='myntra_db',
                user='root',
                password='password',
                connect_timeout=10
            )
        
        return connection
    except Error as e:
        st.error(f"❌ Business Database connection failed: {e}")
        return None
# ================= USER AUTHENTICATION =================
def hash_password(password: str) -> str:
    return hashlib.sha256(password.encode()).hexdigest()
def create_user(username: str, password: str, email: str = None) -> Tuple[bool, str]:
    if not st.session_state.auth_db:
        return False, "Auth database not connected"
    
    cursor = st.session_state.auth_db.cursor()
    try:
        password_hash = hash_password(password)
        cursor.execute(
            "INSERT INTO users (username, password_hash, email) VALUES (%s, %s, %s)",
            (username, password_hash, email)
        )
        st.session_state.auth_db.commit()
        return True, "User created successfully"
    except Error as e:
        if "Duplicate entry" in str(e):
            return False, "Username already exists"
        return False, f"Error creating user: {e}"
    finally:
        cursor.close()
def authenticate_user(username: str, password: str) -> Optional[int]:
    if not st.session_state.auth_db:
        return None
    
    cursor = st.session_state.auth_db.cursor()
    try:
        password_hash = hash_password(password)
        cursor.execute(
            "SELECT user_id FROM users WHERE username = %s AND password_hash = %s",
            (username, password_hash)
        )
        result = cursor.fetchone()
        return result[0] if result else None
    except Error as e:
        st.error(f"Authentication error: {e}")
        return None
    finally:
        cursor.close()
# ================= CHAT MANAGEMENT (IMPROVED) =================
def generate_smart_chat_title(mode: str, first_question: str = None) -> str:
    """Generate a smart title for new chat"""
    if first_question and len(first_question) > 10:
        # Use first few words of the question
        words = first_question.split()[:5]
        title = ' '.join(words)
        if len(first_question) > 30:
            title += "..."
        return title[:50]  # Limit to 50 characters
    else:
        # Generate a descriptive title based on mode and time
        mode_names = {
            MODE_DATABASE: "Database Query",
            MODE_RAG: "Document Q&A"
        }
        timestamp = datetime.now().strftime("%b %d, %I:%M %p")
        return f"{mode_names.get(mode, 'Chat')} - {timestamp}"

def create_new_chat(user_id: int, title: str = None, mode: str = MODE_DATABASE, first_question: str = None) -> Optional[int]:
    """Create new chat with improved title generation"""
    if not st.session_state.auth_db:
        return None
    
    # Generate smart title if not provided
    if not title:
        title = generate_smart_chat_title(mode, first_question)
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute(
            "INSERT INTO chats (user_id, title, mode) VALUES (%s, %s, %s)",
            (user_id, title, mode)
        )
        st.session_state.auth_db.commit()
        return cursor.lastrowid
    except Error as e:
        st.error(f"Error creating chat: {e}")
        return None
    finally:
        cursor.close()

def rename_chat(chat_id: int, user_id: int, new_title: str) -> Tuple[bool, str]:
    """Rename an existing chat"""
    if not st.session_state.auth_db:
        return False, "Auth database not connected"
    
    if not new_title or len(new_title.strip()) == 0:
        return False, "Title cannot be empty"
    
    if len(new_title) > 255:
        return False, "Title too long (max 255 characters)"
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute(
            "UPDATE chats SET title = %s WHERE chat_id = %s AND user_id = %s",
            (new_title.strip(), chat_id, user_id)
        )
        st.session_state.auth_db.commit()
        
        if cursor.rowcount > 0:
            return True, "Chat renamed successfully"
        else:
            return False, "Chat not found or unauthorized"
    except Error as e:
        return False, f"Error renaming chat: {e}"
    finally:
        cursor.close()

def get_user_chats(user_id: int) -> List[Dict]:
    if not st.session_state.auth_db:
        return []
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute(
            "SELECT chat_id, title, mode, created_at FROM chats WHERE user_id = %s ORDER BY created_at DESC",
            (user_id,)
        )
        return [
            {"chat_id": row[0], "title": row[1], "mode": row[2], "created_at": row[3]}
            for row in cursor.fetchall()
        ]
    except Error as e:
        st.error(f"Error fetching chats: {e}")
        return []
    finally:
        cursor.close()
def get_chat_history(chat_id: int, user_id: int) -> List[Dict]:
    if not st.session_state.auth_db:
        return []
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute("""
            SELECT question, query_generated, response, mode, timestamp
            FROM chat_history
            WHERE chat_id = %s AND user_id = %s
            ORDER BY timestamp ASC
        """, (chat_id, user_id))
        
        return [
            {
                "question": row[0],
                "query": row[1],
                "response": row[2],
                "mode": row[3],
                "timestamp": row[4]
            }
            for row in cursor.fetchall()
        ]
    except Error as e:
        st.error(f"Error fetching chat history: {e}")
        return []
    finally:
        cursor.close()
def save_chat_turn(chat_id: int, user_id: int, question: str,
                   query: Optional[str], response: str, mode: str):
    if not st.session_state.auth_db:
        return
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute("""
            INSERT INTO chat_history
            (chat_id, user_id, question, query_generated, response, mode)
            VALUES (%s, %s, %s, %s, %s, %s)
        """, (chat_id, user_id, question, query, response, mode))
        st.session_state.auth_db.commit()
    except Error as e:
        st.error(f"Error saving chat turn: {e}")
    finally:
        cursor.close()
def delete_chat(chat_id: int, user_id: int):
    if not st.session_state.auth_db:
        return
    
    cursor = st.session_state.auth_db.cursor()
    try:
        cursor.execute(
            "DELETE FROM chats WHERE chat_id = %s AND user_id = %s",
            (chat_id, user_id)
        )
        st.session_state.auth_db.commit()
    except Error as e:
        st.error(f"Error deleting chat: {e}")
    finally:
        cursor.close()
# ================= FILE TYPE DETECTION =================
def detect_file_mode(filename: str) -> str:
    """Detect whether file should use database or RAG mode"""
    ext = os.path.splitext(filename)[1].lower()
    
    database_extensions = ['.csv', '.xlsx', '.xls']
    rag_extensions = ['.pdf', '.txt', '.md', '.docx', '.doc', '.pptx', '.ppt',
                     '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif']
    
    if ext in database_extensions:
        return MODE_DATABASE
    elif ext in rag_extensions:
        return MODE_RAG
    else:
        return None
    
    
# ================= CSV/EXCEL TO DATABASE =================
def create_temp_database_from_file(file_bytes: bytes, filename: str) -> Tuple[bool, str, str]:
    """Create temporary database from CSV/Excel file"""
    ext = os.path.splitext(filename)[1].lower()
   
    try:
        if ext == '.csv':
            df = pd.read_csv(BytesIO(file_bytes))
        elif ext in ['.xlsx', '.xls']:
            df = pd.read_excel(BytesIO(file_bytes))
        else:
            return False, None, "Unsupported file format"
       
        # Create table name from filename (unique per user)
        table_name = re.sub(r'[^a-zA-Z0-9_]', '_', os.path.splitext(filename)[0].lower())
        table_name = f"temp_{table_name}_{st.session_state.user_id}"
       
        # Create connection string for temporary database
        if "temp_database" in st.secrets:
            connection = mysql.connector.connect(
                host=st.secrets["temp_database"]["host"],
                port=int(st.secrets["temp_database"]["port"]),
                database=st.secrets["temp_database"]["database"],
                user=st.secrets["temp_database"]["user"],
                password=st.secrets["temp_database"]["password"],
                ssl_disabled=False,
                ssl_verify_cert=True,
                ssl_ca=st.secrets["temp_database"]["ssl_ca"],
                ssl_verify_identity=True
            )
        else:
            connection = mysql.connector.connect(
                host='localhost',
                database='temp_db',
                user='root',
                password='password'
            )
       
        cursor = connection.cursor()
       
        # Drop table if exists (clean slate)
        cursor.execute(f"DROP TABLE IF EXISTS {table_name}")
       
        # Create table schema with AUTO_INCREMENT PRIMARY KEY 'id' column
        create_columns = ["id INT AUTO_INCREMENT PRIMARY KEY"]  # Add this as the first column
        
        # Generate columns from DataFrame (sanitize names, detect types)
        df_columns = []  # We'll build sanitized column names here
        for col in df.columns:
            # Sanitize column name (MySQL-safe: alphanumeric + underscores)
            col_name = re.sub(r'[^a-zA-Z0-9_]', '_', str(col).lower())
            df_columns.append(col_name)  # Track for INSERT later
            
            # Detect column type
            if pd.api.types.is_integer_dtype(df[col]):
                col_type = "INT"
            elif pd.api.types.is_float_dtype(df[col]):
                col_type = "DECIMAL(10,2)"
            else:
                col_type = "TEXT"  # Default to TEXT for safety (handles strings, dates, etc.)
            
            create_columns.append(f"{col_name} {col_type}")
       
        # Build CREATE TABLE query
        create_query = f"CREATE TABLE {table_name} ({', '.join(create_columns)})"
        cursor.execute(create_query)
       
        # Insert data: Specify columns (exclude 'id') and use placeholders for df values only
        if not df.empty:
            columns_str = ", ".join(df_columns)  # e.g., "name, price, category"
            placeholders = ', '.join(['%s'] * len(df_columns))  # e.g., "%s, %s, %s"
            insert_query = f"INSERT INTO {table_name} ({columns_str}) VALUES ({placeholders})"
            
            # Batch insert for efficiency (insert all rows at once)
            data_to_insert = [tuple(row) for _, row in df.iterrows()]  # List of tuples
            cursor.executemany(insert_query, data_to_insert)
       
        connection.commit()
        cursor.close()
       
        # Store connection and table in session state
        st.session_state.temp_db_connection = connection
        st.session_state.temp_table_name = table_name
       
        return True, table_name, f"Created table '{table_name}' with {len(df)} rows and {len(df.columns) + 1} columns (incl. auto-ID)"
       
    except Exception as e:
        # Ensure cursor is closed on error
        try:
            cursor.close()
        except:
            pass
        return False, None, f"Error creating database: {str(e)}"  


# ================= DATABASE SCHEMA =================
def get_database_schema(connection, table_name: str = None) -> Dict:
    """Get database schema"""
    if not connection:
        return {}
    
    schema = {}
    cursor = connection.cursor()
    try:
        if table_name:
            tables = [table_name]
        else:
            cursor.execute("SHOW TABLES")
            tables = [table[0] for table in cursor.fetchall()]
        
        for table in tables:
            cursor.execute(f"DESCRIBE {table}")
            columns = cursor.fetchall()
            schema[table] = {
                "columns": [
                    {"name": col[0], "type": col[1]}
                    for col in columns
                ]
            }
            
            cursor.execute(f"SELECT * FROM {table} LIMIT 2")
            sample_data = cursor.fetchall()
            schema[table]["sample_data"] = sample_data
    except Error as e:
        st.error(f"Error fetching schema: {e}")
    finally:
        cursor.close()
    
    return schema
def format_schema_for_llm(schema: Dict) -> str:
    """Format schema for LLM with emphasis on table names"""
    if not schema:
        return "No schema available"
    
    schema_text = "=== DATABASE SCHEMA ===\n\n"
    
    for table_name, table_info in schema.items():
        schema_text += f"**TABLE NAME: {table_name}** (USE THIS EXACT NAME IN QUERIES)\n"
        schema_text += f"Columns:\n"
        for col in table_info["columns"]:
            schema_text += f"  - {col['name']} ({col['type']})\n"
        
        if table_info.get("sample_data"):
            schema_text += f"\nSample Data (first 2 rows):\n"
            for i, row in enumerate(table_info["sample_data"][:2], 1):
                schema_text += f"  Row {i}: {row}\n"
        
        schema_text += "\n" + "="*50 + "\n\n"
    
    return schema_text
# ================= SEMANTIC MEMORY =================
@st.cache_resource
def load_embedding_model():
    if not SEMANTIC_AVAILABLE:
        return None
    try:
        model = SentenceTransformer('all-MiniLM-L6-v2')
        return model
    except Exception as e:
        return None

class SemanticMemory:
    def __init__(self, embedding_model):
        self.model = embedding_model
        self.messages = []
        self.embeddings = []
    
    def add_message(self, question: str, response: str, query: str = None):
        if not self.model:
            return
        
        text = f"Q: {question}\nA: {response}"
        if query:
            text += f"\nQuery: {query}"
        
        try:
            embedding = self.model.encode(text)
            self.messages.append({"question": question, "response": response, "query": query})
            self.embeddings.append(embedding)
        except:
            pass
    
    def retrieve_relevant(self, query: str, top_k: int = MAX_SEMANTIC_MESSAGES) -> List[Dict]:
        if not self.model or not self.embeddings:
            return []
        
        try:
            query_embedding = self.model.encode(query)
            similarities = np.dot(self.embeddings, query_embedding)
            top_indices = np.argsort(similarities)[-top_k:][::-1]
            
            relevant = []
            for idx in top_indices:
                if similarities[idx] > 0.3:
                    relevant.append(self.messages[idx])
            
            return relevant
        except:
            return []

def build_semantic_memory(chat_history: List[Dict]) -> SemanticMemory:
    embedding_model = load_embedding_model()
    memory = SemanticMemory(embedding_model)
    
    for turn in chat_history:
        memory.add_message(turn["question"], turn["response"], turn.get("query"))
    
    return memory
# ================= FILTER EXTRACTION =================
def extract_active_filters_from_history(chat_history: List[Dict]) -> Tuple[Dict[str, str], str]:
    """Extract active filters from the most recent SQL query in history"""
    active_filters = {
        "category": None,
        "sub_category": None,
        "color": None,
        "size": None,
        "brand": None,
        "price_range": None,
        # Add more as needed
    }
    
    # Check if chat history is empty
    if not chat_history:
        return active_filters, ""
    
    # Get the most recent query
    recent_turns = chat_history[-3:] if len(chat_history) >= 3 else chat_history
    for turn in reversed(recent_turns):  # Start from most recent
        if turn.get("query"):
            query_lower = turn["query"].lower()
            
            # Parse common filters (simple regex patterns)
            if "category" in query_lower and active_filters["category"] is None:
                match = re.search(r"category\s*=\s*'([^']+)'", query_lower)
                if match:
                    active_filters["category"] = match.group(1)
            
            if "sub_category" in query_lower and active_filters["sub_category"] is None:
                match = re.search(r"sub_category\s*=\s*'([^']+)'", query_lower)
                if match:
                    active_filters["sub_category"] = match.group(1)
            
            if "color" in query_lower and active_filters["color"] is None:
                match = re.search(r"color\s*like\s*'%([^']+)%'", query_lower)
                if match:
                    active_filters["color"] = match.group(1)
            
            if "size" in query_lower and active_filters["size"] is None:
                match = re.search(r"size\s*=\s*(\d+)", query_lower)
                if match:
                    active_filters["size"] = match.group(1)
            
            if "brand" in query_lower and active_filters["brand"] is None:
                match = re.search(r"brand\s*like\s*'%([^']+)%'", query_lower)
                if match:
                    active_filters["brand"] = match.group(1)
            
            # Price range example
            if "price" in query_lower and active_filters["price_range"] is None:
                match = re.search(r"price\s*(between|>=|<=)\s*(\d+)", query_lower)
                if match:
                    active_filters["price_range"] = f"price {match.group(1)} {match.group(2)}"
    
    # Build filter string for context only if there are active filters
    has_filters = any(value is not None for value in active_filters.values())
    if not has_filters:
        return active_filters, ""
    
    filter_str = "=== ACTIVE FILTERS FROM PREVIOUS QUERIES ===\n"
    for key, value in active_filters.items():
        if value:
            filter_str += f"- {key}: {value}\n"
    filter_str += "ALWAYS include these filters in the new query unless the user explicitly changes or removes them.\n"
    
    return active_filters, filter_str
# ================= CONTEXT MANAGEMENT =================
def summarize_old_messages(messages: List[Dict]) -> str:
    if not messages:
        return ""
    
    conversation_text = ""
    for msg in messages:
        conversation_text += f"User: {msg['question']}\nAssistant: {msg['response']}\n\n"
    
    system_prompt = """Summarize this conversation focusing on: previous queries, responses, and user intent. Keep under 150 words."""

    user_prompt = f"Summarize:\n\n{conversation_text}\n\nSummary:"

    summary = call_llm(system_prompt, user_prompt, temperature=0.3, max_tokens=300)
    return summary if summary else "Earlier conversation available."

def estimate_tokens(text: str) -> int:
    return len(text) // 4

def build_optimized_context(chat_history: List[Dict], current_question: str) -> Tuple[str, Dict]:
    context_parts = []
    stats = {
        "total_messages": len(chat_history),
        "summarized_count": 0,
        "recent_count": 0,
        "semantic_count": 0,
        "total_tokens": 0
    }
    
    # Extract active filters
    active_filters, filter_section = extract_active_filters_from_history(chat_history)
    if filter_section.strip():
        context_parts.append(filter_section)
        stats["total_tokens"] += estimate_tokens(filter_section)
    
    if len(chat_history) > MAX_RECENT_MESSAGES:
        old_messages = chat_history[:-MAX_RECENT_MESSAGES]
        recent_messages = chat_history[-MAX_RECENT_MESSAGES:]
        
        summary = summarize_old_messages(old_messages)
        if summary:
            context_parts.append(f"=== CONVERSATION SUMMARY ===\n{summary}\n")
            stats["summarized_count"] = len(old_messages)
            stats["total_tokens"] += estimate_tokens(summary)
    else:
        recent_messages = chat_history
    
    if recent_messages:
        context_parts.append("=== RECENT CONVERSATION ===")
        for msg in recent_messages:
            turn = f"User: {msg['question']}\nAssistant: {msg['response']}\n"
            if msg.get('query'):
                turn += f"Query: {msg['query']}\n"
            context_parts.append(turn)
            stats["total_tokens"] += estimate_tokens(turn)
        stats["recent_count"] = len(recent_messages)
    
    if SEMANTIC_AVAILABLE and len(chat_history) > MAX_RECENT_MESSAGES:
        semantic_memory = build_semantic_memory(chat_history[:-MAX_RECENT_MESSAGES])
        relevant_messages = semantic_memory.retrieve_relevant(current_question)
        
        if relevant_messages:
            context_parts.append("\n=== RELATED PAST CONVERSATIONS ===")
            for msg in relevant_messages:
                turn = f"User: {msg['question']}\nAssistant: {msg['response']}\n"
                if msg.get('query'):
                    turn += f"Query: {msg['query']}\n"
                context_parts.append(turn)
                stats["total_tokens"] += estimate_tokens(turn)
            stats["semantic_count"] = len(relevant_messages)
    
    context = "\n".join(context_parts)
    
    if stats["total_tokens"] > MAX_CONTEXT_TOKENS:
        context_parts_reduced = [p for p in context_parts if "RELATED PAST" not in p]
        context = "\n".join(context_parts_reduced)
        stats["semantic_count"] = 0
        stats["total_tokens"] = estimate_tokens(context)
    
    return context, stats
# ================= LLM CALL =================
def call_llm(system_prompt: str, user_prompt: str, temperature: float = 0.1, max_tokens: int = 1000) -> str:
    """Call LLM API"""
    payload = {
        "model": MODEL_NAME,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ],
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    try:
        response = requests.post(API_URL, headers=HEADERS, json=payload, timeout=30)
        if response.status_code == 200:
            return response.json()["choices"][0]["message"]["content"]
        return None
    except Exception as e:
        st.error(f"API failed: {e}")
        return None
# ================= DATABASE MODE FUNCTIONS =================
def is_safe_query(query: str) -> Tuple[bool, str]:
    """Check if SQL query is safe"""
    query_lower = query.lower().strip()
    if not query_lower.startswith('select'):
        return False, "Only SELECT allowed"
    
    # dangerous = ['drop', 'delete', 'truncate', 'alter', 'create', 'insert', 'update']
    dangerous = ['drop', 'delete', 'truncate']
    for keyword in dangerous:
        if keyword in query_lower:
            return False, f"Dangerous keyword: {keyword}"
    
    return True, "Safe"
def generate_sql_query(user_question: str, schema: str, context: str = "") -> Dict:
    """Generate SQL query from natural language"""
    # Extract table names from schema
    table_names = []
    for line in schema.split('\n'):
        if '**TABLE NAME:' in line or line.startswith('Table: '):
            # Extract table name from both old and new formats
            if '**TABLE NAME:' in line:
                table_name = line.split('**TABLE NAME:')[1].split('**')[0].strip()
            else:
                table_name = line.replace('Table: ', '').strip()
            table_names.append(table_name)
    
    # Use the first table name or default to 'footwear_catalog'
    primary_table = table_names[0] if table_names else 'footwear_catalog'
    
    # Enhanced system prompt with few-shot examples for cumulative intent
    few_shot_examples = f"""
EXAMPLE 1:
Previous context: User asked for white sneakers. Last query: SELECT * FROM {primary_table} WHERE LOWER(color) LIKE '%white%' AND sub_category = 'sneakers'
New question: I want size 11
Generated SQL: SELECT * FROM {primary_table} WHERE LOWER(color) LIKE '%white%' AND sub_category = 'sneakers' AND size = 11

EXAMPLE 2:
Previous context: Active filters: color='white', size=11, sub_category='sneakers'
New question: I want Skechers
Generated SQL: SELECT * FROM {primary_table} WHERE LOWER(color) LIKE '%white%' AND sub_category = 'sneakers' AND size = 11 AND LOWER(brand) LIKE '%skechers%'

RULE: ALWAYS build cumulatively. Include ALL active filters from context (e.g., color, size, category) unless the user explicitly says to change/remove them (e.g., 'change color to black')."""
    
    system_prompt = f"""You are a SQL query generator for MySQL database.

{schema}

CRITICAL INSTRUCTIONS:
- The actual table name is: {primary_table}
- ALWAYS use "{primary_table}" in your queries
- NEVER use generic names like "table", "products", "items"

{few_shot_examples}

RULES:
1. ALWAYS use the exact table name "{primary_table}" in FROM clause
2. Use conversation context for cumulative intent - extract and include all active filters from the === ACTIVE FILTERS === section
3. ONLY add/modify filters based on the new question; NEVER drop prior ones unless explicitly contradicted
4. Use LOWER() with LIKE for text matching (e.g., LOWER(color) LIKE '%blue%')
5. Return ONLY the SQL query, no markdown formatting, no explanations, no extra text

REQUIRED FORMAT: SELECT * FROM {primary_table} WHERE ..."""
    
    user_prompt = f"""{context}

QUESTION: {user_question}

Generate a SQL query using table "{primary_table}":"""
    llm_output = call_llm(system_prompt, user_prompt, temperature=0.1, max_tokens=400)
    
    if not llm_output:
        return {"success": False, "error": "LLM failed to generate response", "query": None}
    
    try:
        # Clean the output
        query = llm_output.replace("```sql", "").replace("```", "").strip()
        
        # Remove any text before SELECT
        if "SELECT" in query.upper():
            query = query[query.upper().find("SELECT"):]
        
        # Fix common table name issues
        # Replace generic table names with the actual table name
        generic_names = ['table', 'products', 'items', 'inventory', 'product_table']
        for generic in generic_names:
            # Case-insensitive replacement
            import re
            pattern = r'\bFROM\s+' + re.escape(generic) + r'\b'
            query = re.sub(pattern, f'FROM {primary_table}', query, flags=re.IGNORECASE)
        
        if not query:
            return {"success": False, "error": "Empty query generated", "query": None}
        
        # Verify the correct table name is used
        if 'FROM' in query.upper() and primary_table not in query:
            # Try to fix it by adding proper table name
            query = re.sub(r'FROM\s+\S+', f'FROM {primary_table}', query, flags=re.IGNORECASE, count=1)
        
        is_safe, msg = is_safe_query(query)
        if not is_safe:
            return {"success": False, "error": msg, "query": query}
        
        return {"success": True, "query": query, "error": None}
        
    except Exception as e:
        return {"success": False, "error": str(e), "query": llm_output if 'llm_output' in locals() else None}
def execute_query(connection, query: str) -> Dict:
    """Execute SQL query"""
    if not connection:
        return {"success": False, "error": "No connection", "data": None, "columns": None}
    
    cursor = connection.cursor()
    try:
        cursor.execute(query)
        columns = [desc[0] for desc in cursor.description] if cursor.description else []
        data = cursor.fetchall()
        return {"success": True, "data": data, "columns": columns, "row_count": len(data), "error": None}
    except Error as e:
        return {"success": False, "error": str(e), "data": None, "columns": None}
    finally:
        cursor.close()
def generate_db_response(user_question: str, query: str, result: Dict, context: str = "") -> str:
    """Generate natural language response for database query"""
    if not result["success"]:
        error_msg = result.get("error", "Unknown error")
        return f"I encountered an error while executing the query: {error_msg}"
    
    if not result["data"]:
        return "No results found for your query."
    
    formatted_data = []
    for row in result["data"][:10]:
        item = {col: val for col, val in zip(result["columns"], row)}
        formatted_data.append(item)
    
    system_prompt = "You are a helpful data assistant. Be brief and friendly. Provide a natural response based on the query results."
    
    user_prompt = f"""{context}

QUESTION: {user_question}
RESULTS ({result['row_count']} items found):
{json.dumps(formatted_data, indent=2, default=str)}

Provide a 2-3 sentence natural language response:"""
    
    response = call_llm(system_prompt, user_prompt, temperature=0.3, max_tokens=200)
    return response if response else f"Found {result['row_count']} items matching your query."
# ================= RAG MODE FUNCTIONS =================
def clean_text(text: str) -> str:
    """Clean extracted text"""
    lines = text.split('\n')
    result = '\n'.join(lines)
    result = re.sub(r'\n{3,}', '\n\n', result)
    return result.strip()


def extract_text_from_pdf(pdf_bytes: bytes) -> List[Dict]:
    """Extract text from PDF"""
    if not PYMUPDF_AVAILABLE:
        return []
    
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        pages = []
        
        for i in range(len(doc)):
            page = doc[i]
            text = page.get_text().strip()
            text = clean_text(text)
            
            if text:
                pages.append({"page": i + 1, "text": text})
        
        doc.close()
        return pages
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return []
def semantic_chunking(text: str, chunk_size: int = SEMANTIC_CHUNK_SIZE) -> List[str]:
    """Semantic chunking: splits text at sentence boundaries"""
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_size = 0
    
    for sentence in sentences:
        words = sentence.split()
        sentence_size = len(words)
        
        if current_size + sentence_size > chunk_size and current_chunk:
            chunks.append(' '.join(current_chunk))
            current_chunk = [sentence]
            current_size = sentence_size
        else:
            current_chunk.append(sentence)
            current_size += sentence_size
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks
def process_documents_for_rag(files_data: List[Dict]) -> Tuple[bool, str]:
    """Process documents and store in vector database"""
    if not CHROMADB_AVAILABLE or not SEMANTIC_AVAILABLE:
        return False, "RAG dependencies not available"
    
    embedding_model, _, chroma_client = load_rag_models()
    if not embedding_model or not chroma_client:
        return False, "Failed to load RAG models"
    
    all_chunks = []
    
    for file_data in files_data:
        # Extract text based on file type
        ext = os.path.splitext(file_data['name'])[1].lower()
        
        if ext == '.pdf':
            pages = extract_text_from_pdf(file_data['bytes'])
        else:
            # Add other file type extractors as needed
            continue
        
        # Chunk the text
        for page in pages:
            chunks = semantic_chunking(page['text'])
            for chunk in chunks:
                all_chunks.append({
                    "text": chunk,
                    "page": page['page'],
                    "doc": file_data['name']
                })
    
    if not all_chunks:
        return False, "No content extracted"
    
    # Generate embeddings and store
    texts = [c["text"] for c in all_chunks]
    embeddings = embedding_model.encode(texts, show_progress_bar=False).tolist()
    
    try:
        chroma_client.delete_collection("documents")
    except:
        pass
    
    col = chroma_client.create_collection("documents")
    col.add(
        documents=texts,
        embeddings=embeddings,
        metadatas=[
            {"page": c["page"], "doc": c["doc"]}
            for c in all_chunks
        ],
        ids=[f"c{i}" for i in range(len(texts))]
    )
    
    return True, f"Processed {len(all_chunks)} chunks from {len(files_data)} files"
def rag_query(question: str, context: str = "") -> Tuple[str, List[Dict]]:
    """Query RAG system"""
    if not CHROMADB_AVAILABLE or not SEMANTIC_AVAILABLE:
        return "RAG system not available", []
    
    embedding_model, cross_encoder, chroma_client = load_rag_models()
    if not embedding_model or not chroma_client:
        return "Failed to load RAG models", []
    
    try:
        col = chroma_client.get_collection("documents")
        
        # Retrieve documents
        q_emb = embedding_model.encode([question]).tolist()
        res = col.query(q_emb, n_results=INITIAL_RETRIEVAL_K)
        
        candidates = []
        for doc, metadata, distance in zip(res["documents"][0], res["metadatas"][0], res["distances"][0]):
            candidates.append({
                "text": doc,
                "page": metadata.get("page", 1),
                "doc": metadata.get("doc", "unknown"),
                "score": 1 - distance
            })
        
        if not candidates:
            return "No relevant information found", []
        
        # Rerank if cross-encoder available
        if cross_encoder:
            pairs = [[question, c["text"]] for c in candidates]
            rerank_scores = cross_encoder.predict(pairs)
            for i, score in enumerate(rerank_scores):
                candidates[i]["score"] = float(score)
            candidates.sort(key=lambda x: x["score"], reverse=True)
        
        top_chunks = candidates[:FINAL_CONTEXT_K]
        
        # Enhanced prompting to prevent hallucination
        context_parts = []
        for i, c in enumerate(top_chunks, 1):
            score_info = f"[Relevance: {c.get('score', 0):.2f}]"
            source_info = f"**Source {i}** (Page {c['page']}, {c.get('doc', 'document')}) {score_info}:"
            context_parts.append(f"{source_info}\n{c['text']}\n")
        
        rag_context = "\n".join(context_parts)
        
        system_prompt = """You are an expert document assistant who provides clear, well-structured, and comprehensive answers based on document content. You format answers with appropriate structure including bullet points, paragraphs, and explanations for easy understanding."""
        
        user_prompt = f"""{context}

**INSTRUCTIONS:**
1. Answer the question using ONLY the document content provided below. Do not use any external knowledge or make assumptions.
2. Provide a comprehensive, well-structured answer.
3. Use bullet points for listing multiple items or key points.
4. Use paragraphs for explanations and descriptions.
5. Include specific details, dates, numbers, and examples from the documents.
6. If the answer involves multiple aspects, organize them clearly.
7. Add context and brief explanations to help understanding.
8. Cite the source number (e.g., "Source 1") when referencing specific information.
9. If the answer is not in the documents, clearly state: "Answer not found in the documents." Do not hallucinate or invent information.
**DOCUMENT CONTENT:**
{rag_context}
**QUESTION:**
{question}
**ANSWER:**"""
        
        answer = call_llm(system_prompt, user_prompt, temperature=0.2, max_tokens=500)
        
        sources = [{"doc": c["doc"], "page": c["page"], "score": c["score"]} for c in top_chunks]
        
        return answer if answer else "Could not generate answer", sources
        
    except Exception as e:
        return f"Error: {str(e)}", []
# ================= SESSION STATE INITIALIZATION =================
if "auth_db" not in st.session_state:
    st.session_state.auth_db = get_auth_db_connection()
if "business_db" not in st.session_state:
    st.session_state.business_db = get_business_db_connection()
    if st.session_state.business_db:
        st.session_state.business_schema = get_database_schema(st.session_state.business_db)
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.user_id = None
    st.session_state.username = None
if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = None
if "current_mode" not in st.session_state:
    st.session_state.current_mode = MODE_DATABASE
if "rag_ready" not in st.session_state:
    st.session_state.rag_ready = False
if "temp_db_connection" not in st.session_state:
    st.session_state.temp_db_connection = None
if "temp_table_name" not in st.session_state:
    st.session_state.temp_table_name = None
if "use_system_db" not in st.session_state:
    st.session_state.use_system_db = True
if "show_rename_dialog" not in st.session_state:
    st.session_state.show_rename_dialog = False
if "rename_chat_id" not in st.session_state:
    st.session_state.rename_chat_id = None
# ================= LOGIN UI =================
if not st.session_state.logged_in:
    st.title("🔐 Login / Sign Up")
    
    tab1, tab2 = st.tabs(["Login", "Sign Up"])
    
    with tab1:
        st.subheader("Login")
        login_username = st.text_input("Username", key="login_user")
        login_password = st.text_input("Password", type="password", key="login_pass")
        
        if st.button("Login", use_container_width=True):
            if login_username and login_password:
                user_id = authenticate_user(login_username, login_password)
                if user_id:
                    st.session_state.logged_in = True
                    st.session_state.user_id = user_id
                    st.session_state.username = login_username
                    # Create welcome chat with smart title
                    chat_id = create_new_chat(user_id, None, MODE_DATABASE, "Welcome to AI Assistant")
                    st.session_state.current_chat_id = chat_id
                    st.success(f"Welcome {login_username}!")
                    st.rerun()
                else:
                    st.error("Invalid credentials")
    
    with tab2:
        st.subheader("Sign Up")
        signup_username = st.text_input("Username", key="signup_user")
        signup_email = st.text_input("Email (optional)", key="signup_email")
        signup_password = st.text_input("Password", type="password", key="signup_pass")
        signup_confirm = st.text_input("Confirm Password", type="password", key="signup_confirm")
        
        if st.button("Sign Up", use_container_width=True):
            if signup_username and signup_password:
                if signup_password != signup_confirm:
                    st.error("Passwords don't match")
                elif len(signup_password) < 6:
                    st.error("Password must be 6+ characters")
                else:
                    success, message = create_user(signup_username, signup_password, signup_email)
                    if success:
                        st.success("Account created! Please login.")
                    else:
                        st.error(message)
    st.stop()
# ================= MAIN APP =================
st.title("🎯 Unified AI Assistant")
st.caption(f"Welcome, **{st.session_state.username}** 👋")
# ================= SIDEBAR =================
with st.sidebar:
    st.header("👤 User Menu")
    
    if st.button("🚪 Logout", use_container_width=True):
        st.session_state.logged_in = False
        st.session_state.user_id = None
        st.session_state.username = None
        st.session_state.current_chat_id = None
        st.rerun()
    
    st.divider()
    
    # Mode Selection
    st.subheader("🎛️ Mode Selection")
    mode_options = {
        MODE_DATABASE: "🗄️ Database Query",
        MODE_RAG: "📚 Document Q&A (RAG)"
    }
    
    selected_mode = st.radio(
        "Select Mode:",
        options=list(mode_options.keys()),
        format_func=lambda x: mode_options[x],
        index=0 if st.session_state.current_mode == MODE_DATABASE else 1
    )
    
    if selected_mode != st.session_state.current_mode:
        st.session_state.current_mode = selected_mode
        # Create new chat for new mode with smart title
        chat_id = create_new_chat(
            st.session_state.user_id,
            None,  # Auto-generate title
            selected_mode
        )
        st.session_state.current_chat_id = chat_id
        st.rerun()
    
    st.divider()
    
    # Database Selection (only in database mode)
    if st.session_state.current_mode == MODE_DATABASE:
        st.subheader("🗄️ Database Selection")
        use_system = st.checkbox(
            "Use System Database",
            value=st.session_state.use_system_db,
            help="Use the pre-configured product database"
        )
        
        if use_system != st.session_state.use_system_db:
            st.session_state.use_system_db = use_system
            st.rerun()
        
        if not use_system:
            st.info("💡 Upload CSV/Excel to create your own database")
            uploaded_db_file = st.file_uploader(
                "Upload Database File",
                type=['csv', 'xlsx', 'xls'],
                key="db_uploader"
            )
            
            if uploaded_db_file and st.button("Create Database"):
                with st.spinner("Creating database..."):
                    file_bytes = uploaded_db_file.read()
                    success, table_name, message = create_temp_database_from_file(
                        file_bytes, uploaded_db_file.name
                    )
                    
                    if success:
                        st.success(f"✅ {message}")
                        st.info(f"Table created: {table_name}")
                    else:
                        st.error(f"❌ {message}")
    
    # File Upload (only in RAG mode)
    if st.session_state.current_mode == MODE_RAG:
        st.subheader("📤 Upload Documents")
        uploaded_rag_files = st.file_uploader(
            "Upload Documents for Q&A",
            type=['pdf', 'txt', 'md', 'docx', 'pptx'],
            accept_multiple_files=True,
            key="rag_uploader"
        )
        
        if uploaded_rag_files and st.button("Process Documents"):
            with st.spinner("Processing documents..."):
                files_data = []
                for f in uploaded_rag_files:
                    files_data.append({
                        'name': f.name,
                        'bytes': f.read()
                    })
                
                success, message = process_documents_for_rag(files_data)
                
                if success:
                    st.session_state.rag_ready = True
                    st.success(f"✅ {message}")
                else:
                    st.error(f"❌ {message}")
    
    st.divider()
    st.subheader("💬 Your Chats")
    
    user_chats = get_user_chats(st.session_state.user_id)
    
    if user_chats:
        # Filter chats by current mode
        mode_chats = [c for c in user_chats if c.get("mode", MODE_DATABASE) == st.session_state.current_mode]
        
        if mode_chats:
            # Display chats with better formatting
            for chat in mode_chats:
                col1, col2, col3 = st.columns([6, 2, 2])
                
                with col1:
                    # Truncate long titles
                    display_title = chat['title']
                    if len(display_title) > 30:
                        display_title = display_title[:27] + "..."
                    
                    if st.button(
                        f"{'📌 ' if chat['chat_id'] == st.session_state.current_chat_id else ''}{display_title}",
                        key=f"chat_{chat['chat_id']}",
                        use_container_width=True,
                        type="primary" if chat['chat_id'] == st.session_state.current_chat_id else "secondary"
                    ):
                        st.session_state.current_chat_id = chat['chat_id']
                        st.rerun()
                
                with col2:
                    # Rename button
                    if st.button("✏️", key=f"rename_{chat['chat_id']}", help="Rename chat"):
                        st.session_state.show_rename_dialog = True
                        st.session_state.rename_chat_id = chat['chat_id']
                        st.rerun()
                
                with col3:
                    # Delete button (only for non-active chats or with confirmation)
                    if st.button("🗑️", key=f"del_{chat['chat_id']}", help="Delete chat"):
                        if chat['chat_id'] == st.session_state.current_chat_id:
                            # If deleting current chat, switch to another or create new
                            other_chats = [c for c in mode_chats if c['chat_id'] != chat['chat_id']]
                            if other_chats:
                                st.session_state.current_chat_id = other_chats[0]['chat_id']
                            else:
                                st.session_state.current_chat_id = None
                        
                        delete_chat(chat['chat_id'], st.session_state.user_id)
                        st.rerun()
        else:
            st.info(f"No {mode_options[st.session_state.current_mode]} chats yet")
    else:
        st.info("No chats yet")
    
    # New Chat Button (prominent)
    st.divider()
    if st.button("➕ New Chat", use_container_width=True, type="primary"):
        new_chat_id = create_new_chat(
            st.session_state.user_id,
            None,  # Auto-generate title
            st.session_state.current_mode
        )
        if new_chat_id:
            st.session_state.current_chat_id = new_chat_id
            st.rerun()

# ================= RENAME DIALOG =================
if st.session_state.show_rename_dialog and st.session_state.rename_chat_id:
    @st.dialog("✏️ Rename Chat")
    def rename_dialog():
        # Get current chat title
        user_chats = get_user_chats(st.session_state.user_id)
        current_chat = next((c for c in user_chats if c['chat_id'] == st.session_state.rename_chat_id), None)
        
        if current_chat:
            new_title = st.text_input(
                "New chat title:",
                value=current_chat['title'],
                max_chars=255,
                key="rename_input"
            )
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("💾 Save", use_container_width=True, type="primary"):
                    if new_title and new_title.strip():
                        success, message = rename_chat(
                            st.session_state.rename_chat_id,
                            st.session_state.user_id,
                            new_title
                        )
                        
                        if success:
                            st.success(message)
                            st.session_state.show_rename_dialog = False
                            st.session_state.rename_chat_id = None
                            st.rerun()
                        else:
                            st.error(message)
                    else:
                        st.error("Title cannot be empty")
            
            with col2:
                if st.button("❌ Cancel", use_container_width=True):
                    st.session_state.show_rename_dialog = False
                    st.session_state.rename_chat_id = None
                    st.rerun()
    
    rename_dialog()

# ================= MAIN CHAT INTERFACE =================
st.header("💬 Chat")
# Display mode indicator
mode_emoji = "🗄️" if st.session_state.current_mode == MODE_DATABASE else "📚"
mode_name = "Database Query" if st.session_state.current_mode == MODE_DATABASE else "Document Q&A"
st.info(f"{mode_emoji} **Current Mode:** {mode_name}")
if not st.session_state.current_chat_id:
    st.info("👈 Create a new chat to get started")
else:
    # Get chat history
    chat_history = get_chat_history(st.session_state.current_chat_id, st.session_state.user_id)
    
    # Display chat history
    for turn in chat_history:
        with st.chat_message("user"):
            st.write(turn["question"])
        
        with st.chat_message("assistant"):
            st.write(turn["response"])
            if turn.get("query"):
                with st.expander("🔍 View Query/Details"):
                    st.code(turn["query"], language="sql" if turn.get("mode") == MODE_DATABASE else "text")
    
    # Chat input
    user_question = st.chat_input(f"Ask about {'data' if st.session_state.current_mode == MODE_DATABASE else 'documents'}...")
    
    if user_question:
        # Auto-update chat title if this is the first message
        if len(chat_history) == 0:
            new_title = generate_smart_chat_title(st.session_state.current_mode, user_question)
            rename_chat(st.session_state.current_chat_id, st.session_state.user_id, new_title)
        
        with st.chat_message("user"):
            st.write(user_question)
        
        with st.chat_message("assistant"):
            with st.spinner("Processing..."):
                # Build optimized context from chat history
                context, context_stats = build_optimized_context(chat_history, user_question)
                
                # Show context building stats
                with st.status("Building context...", expanded=False) as status:
                    st.write(f"📊 Context Statistics:")
                    st.write(f"- Total messages in history: {context_stats['total_messages']}")
                    st.write(f"- Summarized messages: {context_stats['summarized_count']}")
                    st.write(f"- Recent messages (full): {context_stats['recent_count']}")
                    st.write(f"- Semantically relevant: {context_stats['semantic_count']}")
                    st.write(f"- Estimated tokens: ~{context_stats['total_tokens']}")
                    status.update(label="Context built!", state="complete")
                
                if st.session_state.current_mode == MODE_DATABASE:
                    # Database mode
                    if st.session_state.use_system_db:
                        if not st.session_state.business_db:
                            response = "⚠️ System database not connected. Please check your database configuration."
                            st.error(response)
                            save_chat_turn(
                                st.session_state.current_chat_id,
                                st.session_state.user_id,
                                user_question,
                                None,
                                response,
                                MODE_DATABASE
                            )
                        else:
                            # Use system database
                            try:
                                schema_text = format_schema_for_llm(st.session_state.business_schema)
                                query_result = generate_sql_query(user_question, schema_text, context)
                                
                                if not query_result["success"]:
                                    response = f"I couldn't generate a valid SQL query. Error: {query_result.get('error', 'Unknown')}"
                                    if query_result.get("query"):
                                        response += f"\n\nAttempted query: {query_result['query']}"
                                    st.error(response)
                                    save_chat_turn(
                                        st.session_state.current_chat_id,
                                        st.session_state.user_id,
                                        user_question,
                                        query_result.get("query"),
                                        response,
                                        MODE_DATABASE
                                    )
                                else:
                                    query = query_result["query"]
                                    result = execute_query(st.session_state.business_db, query)
                                    
                                    if not result["success"]:
                                        response = f"Query execution failed: {result.get('error', 'Unknown error')}\n\nQuery: {query}"
                                        st.error(response)
                                    else:
                                        response = generate_db_response(user_question, query, result, context)
                                        st.write(response)
                                        
                                        with st.expander("🔍 View Query & Results"):
                                            st.subheader("📊 Context Optimization Stats")
                                            col1, col2, col3 = st.columns(3)
                                            with col1:
                                                st.metric("Total Messages", context_stats['total_messages'])
                                                st.metric("Recent (Full)", context_stats['recent_count'])
                                            with col2:
                                                st.metric("Summarized", context_stats['summarized_count'])
                                                st.metric("Semantic", context_stats['semantic_count'])
                                            with col3:
                                                st.metric("Context Tokens", f"~{context_stats['total_tokens']}")
                                                efficiency = ((1 - context_stats['total_tokens'] / 
                                                             max(1, context_stats['total_messages'] * 150)) * 100)
                                                st.metric("Token Efficiency", f"{efficiency:.0f}%")
                                            
                                            st.subheader("🎯 Context Sent to LLM")
                                            st.text_area("Context:", context, height=200, key="context_view_1")
                                            
                                            st.subheader("📝 Generated SQL Query")
                                            st.code(query, language="sql")
                                            
                                            if result["success"] and result["data"]:
                                                st.subheader("📊 Query Results")
                                                df = pd.DataFrame(result["data"][:50], columns=result["columns"])
                                                st.dataframe(df)
                                    
                                    save_chat_turn(
                                        st.session_state.current_chat_id,
                                        st.session_state.user_id,
                                        user_question,
                                        query,
                                        response,
                                        MODE_DATABASE
                                    )
                            except Exception as e:
                                response = f"An unexpected error occurred: {str(e)}"
                                st.error(response)
                                st.exception(e)
                                save_chat_turn(
                                    st.session_state.current_chat_id,
                                    st.session_state.user_id,
                                    user_question,
                                    None,
                                    response,
                                    MODE_DATABASE
                                )
                    else:
                        # Use temporary database
                        if not st.session_state.temp_db_connection:
                            response = "⚠️ Please upload and create a database first"
                            st.error(response)
                            save_chat_turn(
                                st.session_state.current_chat_id,
                                st.session_state.user_id,
                                user_question,
                                None,
                                response,
                                MODE_DATABASE
                            )
                        else:
                            try:
                                schema = get_database_schema(
                                    st.session_state.temp_db_connection,
                                    st.session_state.temp_table_name
                                )
                                schema_text = format_schema_for_llm(schema)
                                query_result = generate_sql_query(user_question, schema_text, context)
                                
                                if not query_result["success"]:
                                    response = f"I couldn't generate a valid SQL query. Error: {query_result.get('error', 'Unknown')}"
                                    if query_result.get("query"):
                                        response += f"\n\nAttempted query: {query_result['query']}"
                                    st.error(response)
                                    save_chat_turn(
                                        st.session_state.current_chat_id,
                                        st.session_state.user_id,
                                        user_question,
                                        query_result.get("query"),
                                        response,
                                        MODE_DATABASE
                                    )
                                else:
                                    query = query_result["query"]
                                    result = execute_query(st.session_state.temp_db_connection, query)
                                    
                                    if not result["success"]:
                                        response = f"Query execution failed: {result.get('error', 'Unknown error')}\n\nQuery: {query}"
                                        st.error(response)
                                    else:
                                        response = generate_db_response(user_question, query, result, context)
                                        st.write(response)
                                        
                                        with st.expander("🔍 View Query & Results"):
                                            st.subheader("📊 Context Optimization Stats")
                                            col1, col2, col3 = st.columns(3)
                                            with col1:
                                                st.metric("Total Messages", context_stats['total_messages'])
                                                st.metric("Recent (Full)", context_stats['recent_count'])
                                            with col2:
                                                st.metric("Summarized", context_stats['summarized_count'])
                                                st.metric("Semantic", context_stats['semantic_count'])
                                            with col3:
                                                st.metric("Context Tokens", f"~{context_stats['total_tokens']}")
                                                efficiency = ((1 - context_stats['total_tokens'] / 
                                                             max(1, context_stats['total_messages'] * 150)) * 100)
                                                st.metric("Token Efficiency", f"{efficiency:.0f}%")
                                            
                                            st.subheader("🎯 Context Sent to LLM")
                                            st.text_area("Context:", context, height=200, key="context_view_2")
                                            
                                            st.subheader("📝 Generated SQL Query")
                                            st.code(query, language="sql")
                                            
                                            if result["success"] and result["data"]:
                                                st.subheader("📊 Query Results")
                                                df = pd.DataFrame(result["data"][:50], columns=result["columns"])
                                                st.dataframe(df)
                                    
                                    save_chat_turn(
                                        st.session_state.current_chat_id,
                                        st.session_state.user_id,
                                        user_question,
                                        query,
                                        response,
                                        MODE_DATABASE
                                    )
                            except Exception as e:
                                response = f"An unexpected error occurred: {str(e)}"
                                st.error(response)
                                st.exception(e)
                                save_chat_turn(
                                    st.session_state.current_chat_id,
                                    st.session_state.user_id,
                                    user_question,
                                    None,
                                    response,
                                    MODE_DATABASE
                                )
                
                else:
                    # RAG mode
                    if not st.session_state.rag_ready:
                        st.error("⚠️ Please upload and process documents first")
                    else:
                        answer, sources = rag_query(user_question, context)
                        st.write(answer)
                        
                        with st.expander("📚 Sources"):
                            for i, source in enumerate(sources, 1):
                                st.write(f"**Source {i}:** {source['doc']} (Page {source['page']}) - Score: {source['score']:.3f}")
                        
                        with st.expander("🔍 View Context & Results"):
                            st.subheader("📊 Context Optimization Stats")
                            col1, col2, col3 = st.columns(3)
                            with col1:
                                st.metric("Total Messages", context_stats['total_messages'])
                                st.metric("Recent (Full)", context_stats['recent_count'])
                            with col2:
                                st.metric("Summarized", context_stats['summarized_count'])
                                st.metric("Semantic", context_stats['semantic_count'])
                            with col3:
                                st.metric("Context Tokens", f"~{context_stats['total_tokens']}")
                                efficiency = ((1 - context_stats['total_tokens'] / 
                                             max(1, context_stats['total_messages'] * 150)) * 100)
                                st.metric("Token Efficiency", f"{efficiency:.0f}%")
                            
                            st.subheader("🎯 Context Sent to LLM")
                            st.text_area("Context:", context, height=200)
                        
                        save_chat_turn(
                            st.session_state.current_chat_id,
                            st.session_state.user_id,
                            user_question,
                            json.dumps(sources),
                            answer,
                            MODE_RAG
                        )
# ================= FOOTER =================
st.divider()
col1, col2, col3 = st.columns(3)
with col1:
    st.caption("🤖 Powered by LLaMA 3.1")
with col2:
    if st.session_state.current_mode == MODE_DATABASE:
        db_status = "System DB" if st.session_state.use_system_db else "Custom DB"
        st.caption(f"🗄️ {db_status}")
    else:
        rag_status = "✅ Ready" if st.session_state.rag_ready else "⚠️ Not Ready"
        st.caption(f"📚 RAG: {rag_status}")
with col3:
    st.caption(f"📊 Mode: {mode_name}")